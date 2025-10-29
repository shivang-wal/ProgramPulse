from fastapi import FastAPI, APIRouter, HTTPException
from fastapi.responses import FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
import uuid
from datetime import datetime, timezone
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile


ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Create the main app without a prefix
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")


class BugSeverity(BaseModel):
    critical: int = 0
    high: int = 0
    medium: int = 0
    low: int = 0

class Project(BaseModel):
    model_config = ConfigDict(extra="ignore")
    
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    status: str  # On Track, At Risk, Delayed, Completed
    completedThisWeek: str = ""
    risks: str = ""
    escalation: str = ""
    plannedNextWeek: str = ""
    bugs: BugSeverity = Field(default_factory=BugSeverity)
    createdAt: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class ProjectCreate(BaseModel):
    name: str
    status: str = "On Track"
    completedThisWeek: str = ""
    risks: str = ""
    escalation: str = ""
    plannedNextWeek: str = ""
    bugs: BugSeverity = Field(default_factory=BugSeverity)

class ProjectUpdate(BaseModel):
    name: Optional[str] = None
    status: Optional[str] = None
    completedThisWeek: Optional[str] = None
    risks: Optional[str] = None
    escalation: Optional[str] = None
    plannedNextWeek: Optional[str] = None
    bugs: Optional[BugSeverity] = None

class CalendarEvent(BaseModel):
    model_config = ConfigDict(extra="ignore")
    
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    date: str  # ISO date string YYYY-MM-DD
    startTime: str = "09:00"  # HH:MM format
    endTime: str = "10:00"   # HH:MM format
    title: str
    description: str = ""
    category: str = "General"  # Event category for grouping
    projectId: Optional[str] = None
    color: str = "#667eea"  # Default purple color
    createdAt: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class CalendarEventCreate(BaseModel):
    date: str
    startTime: str = "09:00"
    endTime: str = "10:00"
    title: str
    description: str = ""
    category: str = "General"
    projectId: Optional[str] = None
    color: str = "#667eea"  # Default purple color

class ProjectHistory(BaseModel):
    model_config = ConfigDict(extra="ignore")
    
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    projectId: str
    projectName: str
    status: str
    completedThisWeek: str = ""
    risks: str = ""
    escalation: str = ""
    plannedNextWeek: str = ""
    bugs: BugSeverity = Field(default_factory=BugSeverity)
    updatedAt: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


# Project Routes
@api_router.post("/projects", response_model=Project)
async def create_project(input: ProjectCreate):
    project_dict = input.model_dump()
    project_obj = Project(**project_dict)
    
    doc = project_obj.model_dump()
    doc['createdAt'] = doc['createdAt'].isoformat()
    
    await db.projects.insert_one(doc)
    return project_obj

@api_router.get("/projects", response_model=List[Project])
async def get_projects():
    projects = await db.projects.find({}, {"_id": 0}).to_list(1000)
    
    for project in projects:
        if isinstance(project.get('createdAt'), str):
            project['createdAt'] = datetime.fromisoformat(project['createdAt'])
    
    return projects

@api_router.get("/projects/{project_id}", response_model=Project)
async def get_project(project_id: str):
    project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if isinstance(project.get('createdAt'), str):
        project['createdAt'] = datetime.fromisoformat(project['createdAt'])
    
    return project

@api_router.put("/projects/{project_id}", response_model=Project)
async def update_project(project_id: str, input: ProjectUpdate):
    update_data = {k: v for k, v in input.model_dump().items() if v is not None}
    
    if not update_data:
        raise HTTPException(status_code=400, detail="No fields to update")
    
    # Get the current project state before updating
    current_project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    
    if not current_project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Save current state to history before updating
    history_entry = ProjectHistory(
        projectId=current_project['id'],
        projectName=current_project['name'],
        status=current_project['status'],
        completedThisWeek=current_project.get('completedThisWeek', ''),
        risks=current_project.get('risks', ''),
        escalation=current_project.get('escalation', ''),
        plannedNextWeek=current_project.get('plannedNextWeek', ''),
        bugs=BugSeverity(**current_project.get('bugs', {})) if current_project.get('bugs') else BugSeverity()
    )
    
    history_doc = history_entry.model_dump()
    history_doc['updatedAt'] = history_doc['updatedAt'].isoformat()
    await db.project_history.insert_one(history_doc)
    
    # Now update the project
    result = await db.projects.update_one(
        {"id": project_id},
        {"$set": update_data}
    )
    
    updated_project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if isinstance(updated_project.get('createdAt'), str):
        updated_project['createdAt'] = datetime.fromisoformat(updated_project['createdAt'])
    
    return updated_project

@api_router.delete("/projects/{project_id}")
async def delete_project(project_id: str):
    result = await db.projects.delete_one({"id": project_id})
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Also delete project history
    await db.project_history.delete_many({"projectId": project_id})
    
    return {"message": "Project deleted successfully"}

@api_router.get("/projects/{project_id}/history", response_model=List[ProjectHistory])
async def get_project_history(project_id: str):
    """Get the update history for a specific project"""
    history = await db.project_history.find(
        {"projectId": project_id}, 
        {"_id": 0}
    ).sort("updatedAt", -1).to_list(1000)
    
    for entry in history:
        if isinstance(entry.get('updatedAt'), str):
            entry['updatedAt'] = datetime.fromisoformat(entry['updatedAt'])
    
    return history


# Calendar Event Routes
@api_router.post("/events", response_model=CalendarEvent)
async def create_event(input: CalendarEventCreate):
    event_dict = input.model_dump()
    event_obj = CalendarEvent(**event_dict)
    
    doc = event_obj.model_dump()
    doc['createdAt'] = doc['createdAt'].isoformat()
    
    await db.events.insert_one(doc)
    return event_obj

@api_router.get("/events", response_model=List[CalendarEvent])
async def get_events():
    events = await db.events.find({}, {"_id": 0}).to_list(1000)
    
    for event in events:
        if isinstance(event.get('createdAt'), str):
            event['createdAt'] = datetime.fromisoformat(event['createdAt'])
    
    return events

@api_router.delete("/events/{event_id}")
async def delete_event(event_id: str):
    result = await db.events.delete_one({"id": event_id})
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Event not found")
    
    return {"message": "Event deleted successfully"}

@api_router.post("/export-ppt")
async def export_projects_ppt(projects: List[Project]):
    """Generate PowerPoint presentation from projects"""
    try:
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define colors (LucyRx theme)
        PURPLE = RGBColor(74, 65, 115)  # #4A4173
        LIGHT_PURPLE = RGBColor(107, 91, 149)  # #6B5B95
        CREAM = RGBColor(255, 249, 230)  # #FFF9E6
        WHITE = RGBColor(255, 255, 255)
        
        # Status colors
        STATUS_COLORS = {
            'On Track': RGBColor(16, 185, 129),
            'At Risk': RGBColor(245, 158, 11),
            'Delayed': RGBColor(239, 68, 68),
            'Completed': RGBColor(99, 102, 241)
        }
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        # Background for title slide
        background = title_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PURPLE
        
        # Add logo
        logo_path = Path(__file__).parent / 'lucy_logo.png'
        if logo_path.exists():
            title_slide.shapes.add_picture(
                str(logo_path),
                Inches(4),
                Inches(1.5),
                width=Inches(2)
            )
        
        # Title
        title_box = title_slide.shapes.add_textbox(
            Inches(0.5), Inches(3.5), Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = 'Program Pulse'
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = title_slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5), Inches(9), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = 'keeping a pulse on all LucyRx initiatives'
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.italic = True
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date and project count
        date_text = datetime.now().strftime('%B %d, %Y')
        info_box = title_slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(9), Inches(0.5)
        )
        info_frame = info_box.text_frame
        info_frame.text = f'{date_text}\n{len(projects)} Active Project{"s" if len(projects) != 1 else ""}'
        for para in info_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        
        # Create slides for each project
        for idx, project in enumerate(projects):
            # Use blank layout
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            # Background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = CREAM
            
            # Header bar
            header_shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0), Inches(0),
                Inches(10), Inches(0.6)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = PURPLE
            header_shape.line.fill.background()
            
            # Add logo to header
            if logo_path.exists():
                slide.shapes.add_picture(
                    str(logo_path),
                    Inches(0.2),
                    Inches(0.15),
                    height=Inches(0.3)
                )
            
            # Project number
            header_text = header_shape.text_frame
            header_text.text = f'Project {idx + 1} of {len(projects)}'
            header_text.paragraphs[0].font.size = Pt(14)
            header_text.paragraphs[0].font.color.rgb = WHITE
            header_text.paragraphs[0].alignment = PP_ALIGN.RIGHT
            header_text.margin_right = Inches(0.3)
            
            # Project name
            name_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1), Inches(9), Inches(0.7)
            )
            name_frame = name_box.text_frame
            name_frame.text = project.name or 'Unnamed Project'
            name_para = name_frame.paragraphs[0]
            name_para.font.size = Pt(32)
            name_para.font.bold = True
            name_para.font.color.rgb = PURPLE
            
            # Status badge
            status_color = STATUS_COLORS.get(project.status, STATUS_COLORS['On Track'])
            status_shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), Inches(1.8),
                Inches(1.5), Inches(0.4)
            )
            status_shape.fill.solid()
            status_shape.fill.fore_color.rgb = status_color
            status_shape.line.fill.background()
            
            status_text = status_shape.text_frame
            status_text.text = project.status or 'On Track'
            status_para = status_text.paragraphs[0]
            status_para.font.size = Pt(14)
            status_para.font.bold = True
            status_para.font.color.rgb = WHITE
            status_para.alignment = PP_ALIGN.CENTER
            
            y_pos = 2.4
            
            # Helper function to add section
            def add_section(title, content, y):
                if not content or content == 'None' or content == 'NA':
                    return y
                
                # Section box
                section_shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(0.5), Inches(y),
                    Inches(9), Inches(0.8)
                )
                section_shape.fill.solid()
                section_shape.fill.fore_color.rgb = RGBColor(245, 240, 255)
                section_shape.line.fill.background()
                
                # Section text
                text_frame = section_shape.text_frame
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_left = Inches(0.2)
                text_frame.margin_right = Inches(0.2)
                
                # Title
                p = text_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = LIGHT_PURPLE
                
                # Content
                p = text_frame.add_paragraph()
                p.text = content
                p.font.size = Pt(12)
                p.font.color.rgb = PURPLE
                p.space_before = Pt(2)
                
                return y + 0.9
            
            # Add sections
            if project.completedThisWeek:
                y_pos = add_section('COMPLETED THIS WEEK', project.completedThisWeek, y_pos)
            
            if project.risks and project.risks != 'None' and project.risks != 'NA':
                y_pos = add_section('RISKS', project.risks, y_pos)
            
            if project.escalation and project.escalation != 'None':
                y_pos = add_section('ESCALATION', project.escalation, y_pos)
            
            if project.plannedNextWeek:
                y_pos = add_section('PLANNED NEXT WEEK', project.plannedNextWeek, y_pos)
            
            # Bug severity matrix
            bugs = project.bugs
            total_bugs = bugs.critical + bugs.high + bugs.medium + bugs.low
            
            if total_bugs > 0 and y_pos < 6.5:
                # Title
                bug_title = slide.shapes.add_textbox(
                    Inches(0.5), Inches(y_pos), Inches(9), Inches(0.3)
                )
                bug_title_frame = bug_title.text_frame
                bug_title_frame.text = f'BUG SEVERITY MATRIX (Total: {total_bugs})'
                bug_title_frame.paragraphs[0].font.size = Pt(11)
                bug_title_frame.paragraphs[0].font.bold = True
                bug_title_frame.paragraphs[0].font.color.rgb = LIGHT_PURPLE
                
                y_pos += 0.35
                
                # Bug cards
                bug_data = [
                    ('CRITICAL', bugs.critical, RGBColor(220, 38, 38)),
                    ('HIGH', bugs.high, RGBColor(245, 158, 11)),
                    ('MEDIUM', bugs.medium, RGBColor(59, 130, 246)),
                    ('LOW', bugs.low, RGBColor(16, 185, 129))
                ]
                
                x_pos = 0.5
                for label, count, color in bug_data:
                    card = slide.shapes.add_shape(
                        1,  # Rectangle
                        Inches(x_pos), Inches(y_pos),
                        Inches(2), Inches(0.6)
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = color
                    card.line.fill.background()
                    
                    card_text = card.text_frame
                    card_text.vertical_anchor = 1  # Middle
                    
                    # Label
                    p = card_text.paragraphs[0]
                    p.text = label
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
                    
                    # Count
                    p = card_text.add_paragraph()
                    p.text = str(count)
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
                    
                    x_pos += 2.2
            
            # Footer
            footer = slide.shapes.add_textbox(
                Inches(0.5), Inches(7), Inches(9), Inches(0.3)
            )
            footer_frame = footer.text_frame
            footer_frame.text = 'Generated by Program Pulse'
            footer_para = footer_frame.paragraphs[0]
            footer_para.font.size = Pt(10)
            footer_para.font.italic = True
            footer_para.font.color.rgb = LIGHT_PURPLE
            footer_para.alignment = PP_ALIGN.CENTER
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        temp_file.close()
        
        # Return file
        filename = f'program_pulse_projects_{datetime.now().strftime("%Y-%m-%d")}.pptx'
        return FileResponse(
            temp_file.name,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename=filename
        )
    
    except Exception as e:
        logging.error(f"Error generating PPT: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate PowerPoint: {str(e)}")


# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()